"""PPTX text and image extractor using python-pptx."""
import io
import hashlib
from pptx import Presentation
from pptx.util import Inches
from google.cloud import storage
import structlog

logger = structlog.get_logger()


class PPTXExtractor:
    """Extracts text content and images from PPTX files."""

    def __init__(self, assets_bucket: str):
        self.storage_client = storage.Client()
        self.assets_bucket = self.storage_client.bucket(assets_bucket)

    def extract(self, pptx_bytes: bytes, lesson_id: str) -> list[dict]:
        """
        Extract slide-level content from a PPTX file.
        Returns a list of dicts with slide_number, title, text_content, image_urls, notes.
        """
        pptx_io = io.BytesIO(pptx_bytes)
        prs = Presentation(pptx_io)
        slides_data = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_data = {
                "slide_number": slide_num,
                "title": "",
                "text_content": "",
                "image_urls": [],
                "notes": "",
            }

            # Extract title
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text.strip()

            # Extract all text content
            text_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_parts.append(text)

                # Extract images
                if shape.shape_type == 13:  # Picture type
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        content_type = image.content_type or "image/png"
                        ext = content_type.split("/")[-1]
                        img_hash = hashlib.md5(image_bytes).hexdigest()[:8]
                        blob_name = f"lessons/{lesson_id}/images/slide_{slide_num}_{img_hash}.{ext}"

                        blob = self.assets_bucket.blob(blob_name)
                        blob.upload_from_string(image_bytes, content_type=content_type)
                        slide_data["image_urls"].append(
                            f"gs://{self.assets_bucket.name}/{blob_name}"
                        )
                    except Exception as e:
                        logger.warning("image_extraction_failed", slide=slide_num, error=str(e))

            slide_data["text_content"] = "\n".join(text_parts)

            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()

            slides_data.append(slide_data)
            logger.info("slide_extracted", slide_number=slide_num, title=slide_data["title"])

        return slides_data
